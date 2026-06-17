from pathlib import Path

import fitz
from docx import Document as DocxDocument
from langchain_core.documents import Document
from pptx import Presentation

from backend.schemas import DocumentPreview


def load_file_to_documents(file_path: str | Path, original_filename: str | None = None) -> list[Document]:
    path = Path(file_path)
    suffix = path.suffix.lower()
    filename = original_filename or path.name

    if suffix == ".pdf":
        return _load_pdf(path, filename)
    if suffix == ".docx":
        return _load_docx(path, filename)
    if suffix == ".pptx":
        return _load_pptx(path, filename)
    if suffix in {".txt", ".md"}:
        return _load_text(path, filename, suffix.lstrip("."))

    raise ValueError(f"unsupported file type for parsing: {suffix or 'unknown'}")


def documents_to_preview(documents: list[Document], max_chars: int = 500) -> list[DocumentPreview]:
    previews = []
    for doc in documents:
        content = doc.page_content
        if len(content) > max_chars:
            content = content[:max_chars] + "..."
        previews.append(DocumentPreview(page_content=content, metadata=dict(doc.metadata)))
    return previews


def _base_metadata(filename: str, file_type: str) -> dict:
    return {
        "source": filename,
        "filename": filename,
        "file_type": file_type,
    }


def _load_pdf(path: Path, filename: str) -> list[Document]:
    documents: list[Document] = []
    with fitz.open(path) as pdf:
        for page_index, page in enumerate(pdf, start=1):
            text = page.get_text("text").strip()
            if not text:
                continue
            metadata = _base_metadata(filename, "pdf")
            metadata["page"] = page_index
            documents.append(Document(page_content=text, metadata=metadata))
    return documents


def _load_docx(path: Path, filename: str) -> list[Document]:
    docx = DocxDocument(path)
    documents: list[Document] = []
    for paragraph_index, paragraph in enumerate(docx.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        metadata = _base_metadata(filename, "docx")
        metadata["paragraph"] = paragraph_index
        documents.append(Document(page_content=text, metadata=metadata))
    return documents


def _load_pptx(path: Path, filename: str) -> list[Document]:
    presentation = Presentation(path)
    documents: list[Document] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_parts.append(text)
        slide_text = "\n".join(text_parts).strip()
        if not slide_text:
            continue
        metadata = _base_metadata(filename, "pptx")
        metadata["slide"] = slide_index
        documents.append(Document(page_content=slide_text, metadata=metadata))
    return documents


def _load_text(path: Path, filename: str, file_type: str) -> list[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    if not text:
        return []
    return [Document(page_content=text, metadata=_base_metadata(filename, file_type))]
